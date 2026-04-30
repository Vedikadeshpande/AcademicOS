"""File processing service — extracts text from PDF/PPT, cleans, chunks, and stores."""
import os
import re
from sqlalchemy import select

from app.database import async_session
from app.models.upload import Upload, ContentChunk
from app.models.syllabus import SyllabusTopic, SyllabusUnit
from app.utils.text_cleaning import clean_text
from app.utils.chunking import chunk_text


async def process_file(upload_id: str):
    """Main background task: extract text, clean, chunk, and store."""
    async with async_session() as db:
        result = await db.execute(select(Upload).where(Upload.id == upload_id))
        upload = result.scalar_one_or_none()
        if not upload:
            return

        try:
            # 1. Extract raw text
            ext = os.path.splitext(upload.file_path)[1].lower()
            if ext == ".pdf":
                pages = extract_pdf_text(upload.file_path)
            elif ext in (".pptx", ".ppt"):
                pages = extract_ppt_text(upload.file_path)
            else:
                upload.status = "error"
                await db.commit()
                return

            # 2. Clean and chunk
            chunk_index = 0
            for page in pages:
                cleaned = clean_text(page["text"])
                if not cleaned or len(cleaned) < 20:
                    continue

                chunks = chunk_text(cleaned)
                for chunk in chunks:
                    content_chunk = ContentChunk(
                        upload_id=upload.id,
                        content=chunk,
                        source_page=str(page.get("page_num", page.get("slide_num", ""))),
                        chunk_index=chunk_index,
                    )
                    db.add(content_chunk)
                    chunk_index += 1

            upload.status = "done"
            await db.commit()

            # 3. Auto-populate content_cache on matching topics
            await _cache_content_for_topics(upload.subject_id, db)

            # 4. Pre-generate quiz questions & flashcards in background (no wait)
            import asyncio
            asyncio.create_task(_run_pregeneration(upload.subject_id))

        except Exception as e:
            upload.status = "error"
            await db.commit()
            raise e


def extract_pdf_text(file_path: str) -> list[dict]:
    """Extract text from PDF using PyPDF2."""
    import PyPDF2

    pages = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"page_num": i + 1, "text": text})
    return pages


def extract_ppt_text(file_path: str) -> list[dict]:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Check if this is the title shape
                if slide.shapes.title and hasattr(slide.shapes.title, 'shape_id'):
                    if shape.shape_id == slide.shapes.title.shape_id:
                        title = text
                        continue
                body_parts.append(text)

        full_text = f"{title}\n{chr(10).join(body_parts)}" if title else "\n".join(body_parts)
        slides.append({
            "slide_num": i + 1,
            "title": title,
            "text": full_text,
        })
    return slides


async def extract_text_from_upload(upload_id: str, db) -> str:
    """Helper: get all extracted text from an upload's chunks."""
    result = await db.execute(
        select(ContentChunk)
        .where(ContentChunk.upload_id == upload_id)
        .order_by(ContentChunk.chunk_index)
    )
    chunks = result.scalars().all()
    if chunks:
        return "\n\n".join(c.content for c in chunks)

    # If no chunks yet, extract directly
    upload_result = await db.execute(select(Upload).where(Upload.id == upload_id))
    upload = upload_result.scalar_one_or_none()
    if not upload:
        return ""

    ext = os.path.splitext(upload.file_path)[1].lower()
    if ext == ".pdf":
        pages = extract_pdf_text(upload.file_path)
    elif ext in (".pptx", ".ppt"):
        pages = extract_ppt_text(upload.file_path)
    else:
        return ""

    return "\n\n".join(clean_text(p["text"]) for p in pages if p["text"])


MAX_CACHE_CHARS = 600  # Keep cache compact per topic


async def _cache_content_for_topics(subject_id: str, db):
    """Match content chunks to syllabus topics and store condensed cache.
    
    This runs once after file upload so quiz/flashcard generation
    can read content_cache directly instead of joining ContentChunk tables.
    """
    # Get all topics for this subject
    topics_result = await db.execute(
        select(SyllabusTopic)
        .join(SyllabusUnit)
        .where(SyllabusUnit.subject_id == subject_id)
    )
    topics = topics_result.scalars().all()
    if not topics:
        return

    # Get all content chunks for this subject
    chunks_result = await db.execute(
        select(ContentChunk.content)
        .join(Upload, ContentChunk.upload_id == Upload.id)
        .where(Upload.subject_id == subject_id)
    )
    all_chunks = [r[0] for r in chunks_result.all() if r[0]]
    if not all_chunks:
        return

    full_content = "\n".join(all_chunks)

    for topic in topics:
        # Extract keywords from topic title (words >= 4 chars)
        keywords = [w.lower() for w in re.split(r'[\s/,;()\-]+', topic.title) if len(w) >= 3]
        if not keywords:
            continue

        # Find sentences that match this topic's keywords
        sentences = re.split(r'(?<=[.!?])\s+', full_content)
        relevant = []
        for sent in sentences:
            sent = sent.strip()
            if len(sent) < 25 or len(sent) > 400:
                continue
            sent_lower = sent.lower()
            # Score: how many keywords match
            matches = sum(1 for kw in keywords if kw in sent_lower)
            if matches > 0:
                relevant.append((matches, sent))

        # Sort by relevance, take top sentences fitting within limit
        relevant.sort(key=lambda x: -x[0])
        cache_parts = []
        total_len = 0
        for _, sent in relevant:
            if total_len + len(sent) > MAX_CACHE_CHARS:
                break
            cache_parts.append(sent)
            total_len += len(sent) + 1

        if cache_parts:
            topic.content_cache = " ".join(cache_parts)

    await db.commit()


async def refresh_content_cache(subject_id: str):
    """Public function to refresh content cache — can be called from routers."""
    async with async_session() as db:
        await _cache_content_for_topics(subject_id, db)


async def _run_pregeneration(subject_id: str):
    """Safely run pre-generation in background (errors don't crash the server)."""
    try:
        from app.services.pregeneration import pregenerate_for_subject
        await pregenerate_for_subject(subject_id)
    except Exception as e:
        print(f"[PreGen] Background pre-generation failed: {e}")


